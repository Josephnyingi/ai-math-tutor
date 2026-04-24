"""
Generate pitch deck for AI Math Tutor (S2.T3.1 AIMS KTT Hackathon).
Run: python3 scripts/make_pitch_deck.py
Output: AI_Math_Tutor_Pitch_Deck.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────
BLUE        = RGBColor(0x1a, 0x56, 0xdb)   # primary brand blue
BLUE_LIGHT  = RGBColor(0xe8, 0xf0, 0xfe)   # light blue bg
DARK        = RGBColor(0x11, 0x18, 0x27)   # near-black
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
AMBER       = RGBColor(0xF5, 0xA6, 0x23)
GREEN       = RGBColor(0x0C, 0xA6, 0x78)
GREY        = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GREY  = RGBColor(0xF3, 0xF4, 0xF6)

W = Inches(13.33)   # widescreen width
H = Inches(7.5)     # widescreen height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── Helper utilities ─────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb


def slide_header(slide, title, subtitle=None, bg=BLUE):
    """Dark header bar at top of slide."""
    add_rect(slide, 0, 0, W, Inches(1.4), fill=bg)
    add_text(slide, title, Inches(0.4), Inches(0.18), Inches(12), Inches(0.7),
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(0.85), Inches(12), Inches(0.45),
                 size=14, bold=False, color=RGBColor(0xc7, 0xd2, 0xfe))


def bullet_block(slide, items, l, t, w, h, size=15, color=DARK, gap=0.42):
    """Render a list of bullet strings."""
    y = t
    for item in items:
        prefix = "•  " if not item.startswith("→") else ""
        add_text(slide, prefix + item, l, y, w, Inches(gap),
                 size=size, color=color, wrap=True)
        y += Inches(gap)


def stat_box(slide, value, label, l, t, w=Inches(2.4), h=Inches(1.5),
             bg=BLUE, vc=WHITE, lc=RGBColor(0xc7, 0xd2, 0xfe)):
    add_rect(slide, l, t, w, h, fill=bg)
    add_text(slide, value, l, t + Inches(0.15), w, Inches(0.75),
             size=30, bold=True, color=vc, align=PP_ALIGN.CENTER)
    add_text(slide, label, l, t + Inches(0.85), w, Inches(0.55),
             size=11, bold=False, color=lc, align=PP_ALIGN.CENTER, wrap=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=DARK)
add_rect(s, 0, 0, Inches(0.18), H, fill=BLUE)          # accent bar

add_text(s, "AI Math Tutor", Inches(0.45), Inches(1.6), Inches(12), Inches(1.4),
         size=54, bold=True, color=WHITE)
add_text(s, "for Early Learners", Inches(0.45), Inches(2.9), Inches(12), Inches(0.8),
         size=36, bold=False, color=RGBColor(0xc7, 0xd2, 0xfe))

add_text(s, "Offline · Adaptive · Multilingual  |  Kinyarwanda · French · English",
         Inches(0.45), Inches(3.9), Inches(10), Inches(0.5),
         size=16, color=GREY)

add_rect(s, Inches(0.45), Inches(4.6), Inches(3.5), Inches(0.05), fill=BLUE)

add_text(s, "Joseph Nyingi   ·   josenyingi@gmail.com",
         Inches(0.45), Inches(4.85), Inches(9), Inches(0.4),
         size=14, color=GREY)
add_text(s, "S2.T3.1  ·  AIMS KTT Hackathon  ·  Tier 3  ·  April 2026",
         Inches(0.45), Inches(5.25), Inches(9), Inches(0.4),
         size=13, color=GREY)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_header(s, "The Problem", "A numeracy crisis hiding in plain sight")

facts = [
    ("50 +", "students per\nteacher in Rwanda"),
    ("0", "cloud tutors that\nwork offline"),
    ("< $50", "tablet — the only\ndevice available"),
    ("3", "languages a child\nmay mix in one answer"),
]
x = Inches(0.4)
for val, lbl in facts:
    stat_box(s, val, lbl, x, Inches(1.6), w=Inches(2.9), h=Inches(1.6))
    x += Inches(3.1)

bullets = [
    "A P1 child who can't count at age 7 struggles with fractions at 9 and algebra at 12.",
    "Parents are often non-literate — they cannot tutor one-to-one.",
    "Every existing digital tutor assumes English, a cloud connection, or a literate adult nearby.",
    "The numeracy gap widens every year because no affordable, offline, multilingual tool exists.",
]
bullet_block(s, bullets, Inches(0.4), Inches(3.5), Inches(12.5), Inches(3.5),
             size=16, color=DARK, gap=0.62)


# ════════════════════════════════════════════════════════════════
# SLIDE 3 — OUR SOLUTION
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_header(s, "Our Solution", "One tutor. Any tablet. Any language. No internet.")

pillars = [
    ("🔇  Fully Offline",
     "Whisper-tiny ASR · BKT knowledge tracing · template + LoRA feedback · AES-256-GCM SQLite\nZero bytes leave the device at inference."),
    ("🌍  Truly Multilingual",
     "Kinyarwanda · French · English · code-switched input.\nDetects mid-sentence language switches and replies in the child's dominant language."),
    ("🧠  Adaptive to Each Child",
     "Bayesian Knowledge Tracing per skill. Selects next item at the child's zone of proximal development.\nDyscalculia early-warning flag after 3 plateau sessions."),
    ("📱  Any Hardware",
     "24 MB on-device footprint (target ≤ 75 MB ✓).\nRuns on a $50 Android tablet. Tap + voice dual input — voice is augmented, not required."),
]
y = Inches(1.55)
for title, body in pillars:
    add_rect(s, Inches(0.3), y, Inches(12.7), Inches(0.88), fill=LIGHT_GREY)
    add_text(s, title, Inches(0.45), y + Inches(0.06), Inches(4), Inches(0.38),
             size=14, bold=True, color=BLUE)
    add_text(s, body, Inches(0.45), y + Inches(0.44), Inches(12.3), Inches(0.42),
             size=12, color=DARK, wrap=True)
    y += Inches(1.02)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 — HOW IT WORKS (Architecture)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_header(s, "How It Works", "7 tasks · 1 offline stack")

components = [
    ("1", "ASR", "Whisper-tiny\n−4.5 semitone child\npitch normalisation"),
    ("2", "Lang Detect", "Lexicon + trigram\nhybrid · EN/FR/KIN\n+ code-switch"),
    ("3", "Knowledge\nTracing", "BKT per skill\nElo baseline\nZPD item selection"),
    ("4", "Visual\nGrounding", "BlobCounter\n100% acc n=1–10\nOWLViT-tiny optional"),
    ("5", "Language\nHead", "TinyLlama 1.1B\nQLoRA rank=8\ntrain_loss = 0.477"),
    ("6", "Privacy\nStore", "AES-256-GCM SQLite\nε-DP sync payload\nPIN learner switch"),
    ("7", "Parent\nReport", "HTML/JSON/text\nIcon-only for\nnon-literate parents"),
]
x = Inches(0.25)
for num, name, detail in components:
    add_rect(s, x, Inches(1.55), Inches(1.7), Inches(3.8), fill=BLUE)
    add_text(s, num, x, Inches(1.6), Inches(1.7), Inches(0.45),
             size=11, color=RGBColor(0xc7, 0xd2, 0xfe), align=PP_ALIGN.CENTER)
    add_text(s, name, x, Inches(2.0), Inches(1.7), Inches(0.65),
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, wrap=True)
    add_text(s, detail, x, Inches(2.75), Inches(1.7), Inches(2.3),
             size=11, color=RGBColor(0xc7, 0xd2, 0xfe), align=PP_ALIGN.CENTER, wrap=True)
    x += Inches(1.82)

add_rect(s, 0, Inches(5.6), W, Inches(0.05), fill=BLUE_LIGHT)
add_text(s, "All 7 components run on-device · 19 / 19 smoke tests pass · "
            "python3 scripts/smoke_test.py",
         Inches(0.4), Inches(5.7), Inches(12.5), Inches(0.4),
         size=12, color=GREY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 — KEY RESULTS
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_header(s, "Key Results", "Every number is seeded & reproducible — python3 scripts/eval_end_to_end.py")

stats = [
    ("0.566", "BKT AUC\n(CI 0.547–0.585)", BLUE),
    ("+0.049", "BKT vs Elo\nAUC delta", GREEN),
    ("100%", "Number-word\nparser accuracy", BLUE),
    ("84.2%", "Language detection\n(4-class F1 ≥ 0.80)", GREEN),
    ("5.92/6", "Feedback rubric\n(100% cases ≥ 5/6)", BLUE),
    ("24 MB", "on-device\nfootprint (< 75 MB)", GREEN),
]
x = Inches(0.3)
row2_y = Inches(4.2)
for i, (val, lbl, col) in enumerate(stats):
    y = Inches(1.55) if i < 3 else row2_y
    xi = Inches(0.3) + (i % 3) * Inches(4.35)
    stat_box(s, val, lbl, xi, y, w=Inches(4.0), h=Inches(1.8), bg=col)

add_text(s,
    "BKT strictly dominates Elo (+0.049 AUC, σ = 5.1 — far outside noise). "
    "Brier/log-loss reported transparently even where Elo wins (calibration vs discrimination trade-off).",
    Inches(0.4), Inches(6.25), Inches(12.5), Inches(0.5),
    size=12, color=GREY, wrap=True)


# ════════════════════════════════════════════════════════════════
# SLIDE 6 — PRODUCT DESIGN
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_header(s, "Product Design", "Built for the real user — not the demo user")

scenarios = [
    ("First 90 Seconds\n(6-year-old, Kinyarwanda)",
     [
         "0 s — Warm KIN voice greeting. Large ▶ button, no reading required.",
         "15 s — Diagnostic probe #1: counting stimulus image + 3 tap buttons.",
         "60 s — 5 probes done. BKT initialised per skill.",
         "90 s — Personalised adaptive session at the child's exact level.",
         "2 silences → tap-only mode. No penalty, no shame language.",
     ]),
    ("Shared Tablet\n(3 children, community centre)",
     [
         "Avatar icons only on home screen — no names displayed.",
         "4-digit PIN shown as coloured dots, not digits.",
         "Each child's data encrypted under a separate AES key.",
         "Per-response SQLite commit — hard reboot loses at most 1 answer.",
         "Works indefinitely without WiFi.",
     ]),
    ("Parent Report\n(non-literate parent)",
     [
         "3 big icons: trend arrow · best skill star · needs-help flag.",
         "5 colour bars — green / amber / red — no numbers or percentages.",
         "Session count shown as dots (●●● = 3 sessions this week).",
         "QR code → 20-second voiced Kinyarwanda summary.",
         "Full report understood in under 60 seconds.",
     ]),
]

x = Inches(0.25)
for title, points in scenarios:
    add_rect(s, x, Inches(1.5), Inches(4.2), Inches(5.7), fill=LIGHT_GREY)
    add_text(s, title, x + Inches(0.1), Inches(1.58), Inches(4.0), Inches(0.65),
             size=13, bold=True, color=BLUE, wrap=True)
    y = Inches(2.35)
    for pt in points:
        add_text(s, "· " + pt, x + Inches(0.1), y, Inches(4.0), Inches(0.5),
                 size=11, color=DARK, wrap=True)
        y += Inches(0.54)
    x += Inches(4.45)


# ════════════════════════════════════════════════════════════════
# SLIDE 7 — WHAT'S ALREADY BUILT
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_header(s, "What's Already Built", "This is not a concept — it is a working system")

done = [
    ("✅", "Full offline pipeline", "Whisper-tiny ASR · BKT · visual grounding · template feedback · SQLite — all running"),
    ("✅", "Real LoRA adapter trained", "TinyLlama-1.1B · QLoRA rank=8 · train_loss=0.477 · 793 s on Tesla T4 · committed in repo"),
    ("✅", "HuggingFace model published", "huggingface.co/Nyingi101/math-tutor-tinyllama-lora — public, downloadable"),
    ("✅", "87-item curriculum", "5 skills · 4 age bands · difficulty 1–10 · EN/FR/KIN stems · reproducible from seed"),
    ("✅", "Reproducible evaluation", "eval_end_to_end.py generates all metrics · figures/*.json committed · CI with 95% bounds"),
    ("✅", "Real interaction logging", "log_interaction() + export_interactions_for_finetuning() — ready to collect real child data"),
    ("✅", "GitHub repo", "github.com/Josephnyingi/ai-math-tutor · all code, weights, notebooks, eval artifacts"),
]

y = Inches(1.55)
for icon, title, detail in done:
    add_rect(s, Inches(0.3), y, Inches(12.7), Inches(0.66), fill=LIGHT_GREY)
    add_text(s, icon + "  " + title,
             Inches(0.45), y + Inches(0.08), Inches(3.8), Inches(0.5),
             size=13, bold=True, color=GREEN)
    add_text(s, detail,
             Inches(4.4), y + Inches(0.1), Inches(8.4), Inches(0.46),
             size=12, color=DARK, wrap=True)
    y += Inches(0.74)


# ════════════════════════════════════════════════════════════════
# SLIDE 8 — LIMITATIONS (HONEST)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_header(s, "What We Haven't Solved Yet", "Honest limitations — and the funded plan to fix them", bg=DARK)

rows = [
    ("Synthetic training data",
     "The LoRA was trained on 2,000 generated pairs — not real child speech.",
     "log_interaction() already captures real data from day 1 of deployment.\nSchool pilot → 500 real interactions → retrain."),
    ("Whisper-tiny KIN accuracy",
     "WER ~40% on Kinyarwanda out-of-the-box. Child pitch makes it harder.",
     "Fine-tune Whisper-small on Mozilla Common Voice KIN + pitch-shifted audio.\nTarget: WER < 15%."),
    ("No learning-gain evidence",
     "System is proven technically. Educational impact is unvalidated.",
     "School pilot: pre/post numeracy assessment · 60 children · 2 classes.\nTarget: +0.15 BKT mastery delta."),
    ("No GGUF on-device yet",
     "TinyLlama adapter not yet merged + quantised to run on $50 tablet.",
     "merge-only script ready. 1 day of work post-pilot. Budget allocated."),
]

y = Inches(1.55)
for limitation, problem, solution in rows:
    add_rect(s, Inches(0.3), y, Inches(12.7), Inches(1.2), fill=LIGHT_GREY)
    add_text(s, limitation, Inches(0.45), y + Inches(0.08), Inches(3.0), Inches(0.4),
             size=13, bold=True, color=RGBColor(0xB9, 0x1C, 0x1C))
    add_text(s, "⚠  " + problem, Inches(3.6), y + Inches(0.08), Inches(4.4), Inches(0.45),
             size=12, color=DARK, wrap=True)
    add_text(s, "→  " + solution, Inches(8.2), y + Inches(0.08), Inches(4.8), Inches(1.0),
             size=11, color=GREEN, wrap=True)
    y += Inches(1.28)


# ════════════════════════════════════════════════════════════════
# SLIDE 9 — THE ASK
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=DARK)
add_rect(s, 0, 0, Inches(0.18), H, fill=BLUE)

add_text(s, "Give me $30,000.", Inches(0.45), Inches(0.5), Inches(12), Inches(1.0),
         size=44, bold=True, color=WHITE)
add_text(s,
    "By Month 3, I will put a working, evidence-backed numeracy tutor\n"
    "in the hands of 60 Rwandan children — with peer-reviewed learning-gain data.",
    Inches(0.45), Inches(1.5), Inches(12.5), Inches(1.0),
    size=18, color=RGBColor(0xc7, 0xd2, 0xfe), wrap=True)

budget = [
    ("2 engineers × 3 months", "$21,000"),
    ("KIN linguistic consultant", "$1,500"),
    ("School pilot — 2 sites, tablets, logistics", "$2,000"),
    ("GPU compute — 4 retraining runs", "$500"),
    ("Teacher training workshops", "$1,000"),
    ("Contingency 15%", "$3,900"),
]
y = Inches(2.7)
for item, cost in budget:
    add_rect(s, Inches(0.4), y, Inches(6.2), Inches(0.38), fill=RGBColor(0x1f, 0x29, 0x3a))
    add_text(s, item, Inches(0.55), y + Inches(0.04), Inches(4.5), Inches(0.32),
             size=12, color=WHITE)
    add_text(s, cost, Inches(5.5), y + Inches(0.04), Inches(1.0), Inches(0.32),
             size=12, bold=True, color=AMBER, align=PP_ALIGN.RIGHT)
    y += Inches(0.44)

add_rect(s, Inches(0.4), y + Inches(0.05), Inches(6.2), Inches(0.44), fill=BLUE)
add_text(s, "Total ask", Inches(0.55), y + Inches(0.09), Inches(4.5), Inches(0.35),
         size=14, bold=True, color=WHITE)
add_text(s, "$29,900", Inches(5.0), y + Inches(0.09), Inches(1.5), Inches(0.35),
         size=14, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

deliverables = [
    "KIN-fine-tuned ASR  —  WER < 15%",
    "Real-data retrained LoRA  —  500+ child interactions",
    "GGUF on $50 tablet  —  fully quantised",
    "Learning-gain evidence  —  60 children, pre/post assessment",
    "Teacher dashboard v1  —  class-level BKT heatmap",
    "Research brief  —  ready for AIMS grant / journal",
]
add_text(s, "Month 3 deliverables:", Inches(7.1), Inches(2.7), Inches(5.8), Inches(0.4),
         size=14, bold=True, color=AMBER)
y2 = Inches(3.2)
for d in deliverables:
    add_text(s, "✓  " + d, Inches(7.1), y2, Inches(5.8), Inches(0.4),
             size=12, color=WHITE, wrap=True)
    y2 += Inches(0.44)


# ════════════════════════════════════════════════════════════════
# SLIDE 10 — 90-DAY PLAN
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=WHITE)
slide_header(s, "90-Day Plan", "Milestone by milestone")

months = [
    ("Month 1", "Deploy & Collect", BLUE, [
        "Deploy to School A — 30 children, P1–P3",
        "Pre-assessment baseline (numeracy per skill)",
        "Collect 500 real interaction triples on-device",
        "GGUF quantise TinyLlama → verify on $50 tablet",
        "Weekly teacher check-in; tune silence thresholds",
    ]),
    ("Month 2", "Retrain & Validate", GREEN, [
        "Export pseudonymised interactions → retrain LoRA",
        "Fine-tune Whisper-small on KIN Common Voice",
        "Ship teacher dashboard v1 (class skill heatmap)",
        "Mid-point assessment at Week 6",
        "Onboard School B — 30 more children",
    ]),
    ("Month 3", "Evidence & Scale", AMBER, [
        "Post-assessment vs baseline + control group",
        "Analyse learning gains (target: +0.15 BKT delta)",
        "Draft research brief (AIMS grant / journal)",
        "Package v1.0 with reproducible training pipeline",
        "Present results · identify School C",
    ]),
]

x = Inches(0.25)
for label, title, col, points in months:
    add_rect(s, x, Inches(1.5), Inches(4.2), Inches(0.55), fill=col)
    add_text(s, label + " — " + title, x + Inches(0.1), Inches(1.56), Inches(4.0), Inches(0.44),
             size=14, bold=True, color=WHITE)
    y = Inches(2.15)
    for pt in points:
        add_rect(s, x, y, Inches(4.2), Inches(0.85), fill=LIGHT_GREY)
        add_text(s, pt, x + Inches(0.12), y + Inches(0.1), Inches(3.95), Inches(0.7),
                 size=12, color=DARK, wrap=True)
        y += Inches(0.92)
    x += Inches(4.45)


# ════════════════════════════════════════════════════════════════
# SLIDE 11 — CLOSING / CONTACT
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, fill=DARK)
add_rect(s, 0, 0, Inches(0.18), H, fill=BLUE)

add_text(s,
    "A child who cannot count at 7\nstruggles with fractions at 9\nand algebra at 12.",
    Inches(0.55), Inches(0.7), Inches(12), Inches(2.2),
    size=34, bold=True, color=WHITE, wrap=True)

add_text(s,
    "We built the tutor that sits with her one-to-one,\n"
    "in her language, on the device she already has,\n"
    "for as long as she needs.",
    Inches(0.55), Inches(3.1), Inches(12), Inches(1.6),
    size=20, color=RGBColor(0xc7, 0xd2, 0xfe), wrap=True)

add_rect(s, Inches(0.55), Inches(4.9), Inches(5.5), Inches(0.04), fill=BLUE)

contact = [
    ("Joseph Nyingi", 22, True, WHITE),
    ("josenyingi@gmail.com", 16, False, RGBColor(0xc7, 0xd2, 0xfe)),
    ("github.com/Josephnyingi/ai-math-tutor", 14, False, GREY),
    ("huggingface.co/Nyingi101/math-tutor-tinyllama-lora", 14, False, GREY),
]
y = Inches(5.1)
for text, sz, bold, col in contact:
    add_text(s, text, Inches(0.55), y, Inches(9), Inches(0.48),
             size=sz, bold=bold, color=col)
    y += Inches(0.48)


# ── Save ──────────────────────────────────────────────────────
OUT = "AI_Math_Tutor_Pitch_Deck.pptx"
prs.save(OUT)
print(f"Saved → {OUT}  ({len(prs.slides)} slides)")
